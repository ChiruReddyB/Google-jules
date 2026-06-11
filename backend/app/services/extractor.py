import os
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import io

class TextExtractor:
    @staticmethod
    def extract_from_pdf(file_content: bytes) -> str:
        pdf_reader = PdfReader(io.BytesIO(file_content))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text

    @staticmethod
    def extract_from_docx(file_content: bytes) -> str:
        doc = DocxDocument(io.BytesIO(file_content))
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text

    @staticmethod
    def extract_from_pptx(file_content: bytes) -> str:
        prs = Presentation(io.BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    @staticmethod
    def extract_from_txt(file_content: bytes) -> str:
        return file_content.decode("utf-8")

    @classmethod
    def extract_text(cls, file_content: bytes, filename: str) -> str:
        extension = filename.split(".")[-1].lower()
        if extension == "pdf":
            return cls.extract_from_pdf(file_content)
        elif extension in ["doc", "docx"]:
            return cls.extract_from_docx(file_content)
        elif extension in ["ppt", "pptx"]:
            return cls.extract_from_pptx(file_content)
        elif extension == "txt":
            return cls.extract_from_txt(file_content)
        else:
            raise ValueError(f"Unsupported file extension: {extension}")
